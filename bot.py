import os
import logging
import asyncio
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import Application, CommandHandler, MessageHandler, CallbackQueryHandler, filters, ContextTypes
from pdf2docx import Converter
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation

logging.basicConfig(format='%(levelname)s - %(message)s', level=logging.INFO)

BOT_TOKEN = "8360235283:AAEVm9ujtALyLhldZfmBk9eATYE9uBFZ7Dw"
# Webhook configuration for 24/7 uptime

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "Salom! Hujjat konvertor bot.\n\n"
        "📄 DOCX ↔ PDF\n"
        "📄 PDF → DOCX\n"
        "📽 PPTX → PDF/DOCX\n\n"
        "Hujjat yuboring!"
    )

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    document = update.message.document
    file_name = document.file_name.lower()
    file_size_mb = document.file_size / (1024 * 1024)
    
    if file_size_mb > 20:
        await update.message.reply_text(f"❌ Fayl katta: {file_size_mb:.1f} MB (max 20 MB)")
        return
    
    try:
        file = await context.bot.get_file(document.file_id)
        ext = os.path.splitext(document.file_name)[1]
        name = f"f_{abs(hash(document.file_id)) % 100000}{ext}"
        path = os.path.abspath(name)
        
        msg = await update.message.reply_text("⏳ Yuklanmoqda...")
        await file.download_to_drive(path)
        await msg.delete()
        
        fid = str(abs(hash(document.file_id)) % 1000000)
        context.user_data[fid] = path
        
        kb = None
        if ext == '.docx':
            kb = [[InlineKeyboardButton("PDF", callback_data=f"docx_pdf:{fid}")]]
        elif ext == '.pdf':
            kb = [[InlineKeyboardButton("DOCX", callback_data=f"pdf_docx:{fid}")]]
        elif ext in ['.pptx', '.ppt']:
            kb = [
                [InlineKeyboardButton("PDF", callback_data=f"pptx_pdf:{fid}")],
                [InlineKeyboardButton("DOCX", callback_data=f"pptx_docx:{fid}")]
            ]
        
        if kb:
            await update.message.reply_text("Format tanlang:", reply_markup=InlineKeyboardMarkup(kb))
        else:
            await update.message.reply_text("❌ Format qo'llab-quvvatlanmaydi")
            if os.path.exists(path):
                os.remove(path)
                
    except Exception as e:
        await update.message.reply_text(f"❌ Xatolik: {str(e)}")

async def button_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    
    action, fid = query.data.split(':', 1)
    path = context.user_data.get(fid)
    
    if not path or not os.path.exists(path):
        await query.message.reply_text("❌ Fayl topilmadi")
        return
    
    try:
        msg = await query.message.reply_text("⏳ Konvertatsiya...")
        
        output = None
        if action == "docx_pdf":
            output = await docx_to_pdf(path)
        elif action == "pdf_docx":
            output = await pdf_to_docx(path)
        elif action == "pptx_pdf":
            output = await pptx_to_pdf(path)
        elif action == "pptx_docx":
            output = await pptx_to_docx(path)
        
        await msg.delete()
        
        if output and os.path.exists(output):
            with open(output, 'rb') as f:
                await query.message.reply_document(document=f, filename=os.path.basename(output))
            await query.message.reply_text("✅ Tayyor!")
            os.remove(output)
        else:
            await query.message.reply_text("❌ Xatolik")
            
    except Exception as e:
        await query.message.reply_text(f"❌ {str(e)}")

async def docx_to_pdf(path):
    out = 'sardor.pdf'
    doc = Document(path)
    pdf = SimpleDocTemplate(out, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    for p in doc.paragraphs:
        if p.text.strip():
            story.append(Paragraph(p.text.replace('<', '&lt;').replace('>', '&gt;'), styles['Normal']))
            story.append(Spacer(1, 6))
    
    pdf.build(story)
    return out

async def pdf_to_docx(path):
    out = 'sardor.docx'
    cv = Converter(path)
    cv.convert(out)
    cv.close()
    return out

async def pptx_to_pdf(path):
    out = 'sardor.pdf'
    prs = Presentation(path)
    pdf = SimpleDocTemplate(out, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    for i, slide in enumerate(prs.slides, 1):
        story.append(Paragraph(f"<b>Slide {i}</b>", styles['Heading1']))
        story.append(Spacer(1, 12))
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                story.append(Paragraph(shape.text.replace('<', '&lt;').replace('>', '&gt;'), styles['Normal']))
                story.append(Spacer(1, 6))
        
        story.append(Spacer(1, 20))
    
    pdf.build(story)
    return out

async def pptx_to_docx(path):
    out = 'sardor.docx'
    prs = Presentation(path)
    doc = Document()
    
    for i, slide in enumerate(prs.slides, 1):
        doc.add_heading(f'Slide {i}', level=1)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                doc.add_paragraph(shape.text)
        doc.add_paragraph()
    
    doc.save(out)
    return out

def main():
    app = Application.builder().token(BOT_TOKEN).build()
    app.add_handler(CommandHandler("start", start))
    app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    app.add_handler(CallbackQueryHandler(button_callback))
    
    # Keep-alive: har 5 daqiqada o'ziga ping yuborish
    async def keep_alive(application):
        while True:
            try:
                await asyncio.sleep(300)  # 5 daqiqa
                await application.bot.get_me()
                print("✅ Keep-alive ping yuborildi")
            except Exception as e:
                print(f"⚠️ Keep-alive xatolik: {e}")
    
    app.post_init = lambda app: asyncio.create_task(keep_alive(app))
    
    print("✅ Bot ishga tushdi! (Keep-alive faol)")
    app.run_polling(allowed_updates=Update.ALL_TYPES)

if __name__ == '__main__':
    main()
